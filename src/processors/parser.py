import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Optional

from ..models.file_record import FileRecord
from ..utils.logger import setup_logger

logger = setup_logger()


class DocumentParser:
    TEXT_EXTS = {"txt", "md", "csv", "json", "html", "log", "py", "js", "ts", "yaml", "yml", "xml"}
    DOC_EXTS = {"docx", "doc"}
    XLS_EXTS = {"xlsx", "xls"}
    PPT_EXTS = {"pptx", "ppt"}
    ARCHIVE_EXTS = {"zip", "rar", "7z"}

    def __init__(self, max_text_length: int = 8000):
        self.max_text_length = max_text_length

    def process(self, record: FileRecord) -> FileRecord:
        text = self._parse(record)
        record.text_content = (text or "")[: self.max_text_length]
        record.log_step("parse", f"提取文本 {len(record.text_content)} 字符")
        return record

    def _parse(self, record: FileRecord) -> str:
        ext = record.file_ext
        path = record.source_path
        try:
            if ext in self.TEXT_EXTS:
                return self._parse_text(path)
            if ext == "pdf":
                return self._parse_pdf(path)
            if ext == "docx":
                return self._parse_docx(path)
            if ext == "doc":
                return self._parse_doc(path)
            if ext == "xlsx":
                return self._parse_xlsx(path)
            if ext == "xls":
                return self._parse_xls(path)
            if ext == "pptx":
                return self._parse_pptx(path)
            if ext == "zip":
                return self._parse_zip(record)
            if ext in {"jpg", "jpeg", "png", "gif", "webp", "bmp"}:
                return f"[图片文件] {record.file_name}，大小 {record.file_size} 字节"
            return ""
        except Exception as e:
            logger.debug(f"解析失败 {record.file_name}: {e}")
            return ""

    def _parse_text(self, path: str) -> str:
        for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
            try:
                with open(path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return ""

    def _parse_pdf(self, path: str) -> str:
        try:
            import pdfplumber
            texts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages[:30]:
                    t = page.extract_text() or ""
                    if t:
                        texts.append(t)
            return "\n".join(texts)
        except Exception as e:
            logger.debug(f"pdfplumber 失败: {e}")
            return ""

    def _parse_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
            return "\n".join(parts)
        except Exception as e:
            logger.debug(f"docx 解析失败: {e}")
            return ""

    def _parse_doc(self, path: str) -> str:
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", path],
                capture_output=True, text=True, timeout=30
            )
            return result.stdout if result.returncode == 0 else ""
        except Exception:
            return ""

    def _parse_xlsx(self, path: str) -> str:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets[:5]:
                parts.append(f"### 工作表: {sheet.title}")
                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    if i > 100:
                        break
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
        except Exception as e:
            logger.debug(f"xlsx 解析失败: {e}")
            return ""

    def _parse_xls(self, path: str) -> str:
        try:
            import xlrd
            wb = xlrd.open_workbook(path)
            parts = []
            for sheet in wb.sheets()[:5]:
                parts.append(f"### 工作表: {sheet.name}")
                for i in range(min(sheet.nrows, 100)):
                    cells = [str(sheet.cell_value(i, j)) for j in range(sheet.ncols)]
                    cells = [c for c in cells if c.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
        except Exception:
            return ""

    def _parse_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    parts.append(f"### Slide {i+1}")
                    parts.extend(texts)
            return "\n".join(parts)
        except Exception as e:
            logger.debug(f"pptx 解析失败（可能兼容性问题）: {e}")
            return ""

    def _parse_zip(self, record: FileRecord) -> str:
        parts = [f"[压缩包] {record.file_name}，包含以下文件："]
        children = []
        try:
            with zipfile.ZipFile(record.source_path, "r") as zf:
                for info in zf.infolist():
                    if info.is_dir():
                        continue
                    name = info.filename
                    if name.startswith("__MACOSX") or name.startswith("."):
                        continue
                    parts.append(f"- {name} ({info.file_size} bytes)")
                    children.append(name)
            record.archive_children = children
        except Exception as e:
            parts.append(f"解压失败: {e}")
        return "\n".join(parts)
